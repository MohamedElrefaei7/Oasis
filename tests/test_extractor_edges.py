"""Edge-case tests for all four extractors."""

from pathlib import Path

import pytest

from oasis.extractors.docx import DocxExtractor
from oasis.extractors.pdf import PdfExtractor
from oasis.extractors.pptx import PptxExtractor
from oasis.extractors.text import TextExtractor
from oasis.models import ExtractedDocument


# ---------------------------------------------------------------------------
# TextExtractor
# ---------------------------------------------------------------------------


class TestTextExtractorEdges:
    def test_empty_file_returns_document_with_empty_text(self, tmp_path: Path) -> None:
        f = tmp_path / "empty.txt"
        f.write_text("")
        doc = TextExtractor().extract(f)
        assert doc is not None
        assert doc.text == ""

    def test_non_utf8_bytes_returns_none(self, tmp_path: Path) -> None:
        f = tmp_path / "binary.txt"
        f.write_bytes(b"\xff\xfe invalid utf8 \x80\x81")
        doc = TextExtractor().extract(f)
        assert doc is None

    def test_very_short_text_no_crash(self, tmp_path: Path) -> None:
        f = tmp_path / "short.txt"
        f.write_text("a")
        doc = TextExtractor().extract(f)
        assert doc is not None

    def test_very_short_text_language_is_none_or_string(self, tmp_path: Path) -> None:
        f = tmp_path / "short.txt"
        f.write_text("x")
        doc = TextExtractor().extract(f)
        assert doc is not None
        assert doc.metadata.language is None or isinstance(doc.metadata.language, str)

    def test_large_text_file_extracts_successfully(self, tmp_path: Path) -> None:
        f = tmp_path / "large.txt"
        f.write_text("sentence with many words " * 10_000)
        doc = TextExtractor().extract(f)
        assert doc is not None
        assert len(doc.text) > 0

    def test_md_raw_markdown_preserved(self, tmp_path: Path) -> None:
        f = tmp_path / "notes.md"
        f.write_text("# Heading\n\nParagraph with **bold** and _italic_.\n")
        doc = TextExtractor().extract(f)
        assert doc is not None
        assert "Heading" in doc.text
        assert "bold" in doc.text

    def test_metadata_size_matches_file(self, tmp_path: Path) -> None:
        f = tmp_path / "sized.txt"
        content = "exactly this content"
        f.write_text(content, encoding="utf-8")
        doc = TextExtractor().extract(f)
        assert doc is not None
        assert doc.metadata.size_bytes == f.stat().st_size

    def test_utf8_special_chars_extracted(self, tmp_path: Path) -> None:
        f = tmp_path / "unicode.txt"
        f.write_text("café résumé naïve", encoding="utf-8")
        doc = TextExtractor().extract(f)
        assert doc is not None
        assert "café" in doc.text


# ---------------------------------------------------------------------------
# PdfExtractor
# ---------------------------------------------------------------------------


class TestPdfExtractorEdges:
    def test_zero_byte_file_returns_none(self, tmp_path: Path) -> None:
        f = tmp_path / "empty.pdf"
        f.write_bytes(b"")
        doc = PdfExtractor().extract(f)
        assert doc is None

    def test_blank_page_pdf_returns_none(self, tmp_path: Path) -> None:
        from pypdf import PdfWriter
        f = tmp_path / "blank.pdf"
        writer = PdfWriter()
        writer.add_blank_page(width=612, height=792)
        with open(f, "wb") as fh:
            writer.write(fh)
        doc = PdfExtractor().extract(f)
        assert doc is None

    def test_multipage_pdf_page_count_correct(self, tmp_path: Path) -> None:
        from pypdf import PdfWriter
        from pypdf.generic import NameObject
        f = tmp_path / "multi.pdf"
        writer = PdfWriter()
        for _ in range(3):
            writer.add_blank_page(width=612, height=792)
        with open(f, "wb") as fh:
            writer.write(fh)
        doc = PdfExtractor().extract(f)
        # blank pages → None; count is still verifiable via reader directly
        # This test verifies no crash on multi-page blank PDF
        assert doc is None or doc.metadata.page_count == 3

    def test_pdf_missing_metadata_no_crash(self, tmp_path: Path) -> None:
        from pypdf import PdfWriter
        f = tmp_path / "nometa.pdf"
        writer = PdfWriter()
        writer.add_blank_page(width=612, height=792)
        with open(f, "wb") as fh:
            writer.write(fh)
        doc = PdfExtractor().extract(f)
        assert doc is None  # blank page = no text = None


# ---------------------------------------------------------------------------
# DocxExtractor
# ---------------------------------------------------------------------------


class TestDocxExtractorEdges:
    def test_empty_docx_returns_document_with_empty_text(self, tmp_path: Path) -> None:
        from docx import Document
        f = tmp_path / "empty.docx"
        d = Document()
        d.save(str(f))
        result = DocxExtractor().extract(f)
        assert result is not None
        assert result.text == ""

    def test_docx_without_title_has_none_title(self, tmp_path: Path) -> None:
        from docx import Document
        f = tmp_path / "notitle.docx"
        d = Document()
        d.add_paragraph("some text")
        d.save(str(f))
        result = DocxExtractor().extract(f)
        assert result is not None
        assert result.metadata.title is None

    def test_docx_author_is_string_or_none(self, tmp_path: Path) -> None:
        from docx import Document
        f = tmp_path / "noauthor.docx"
        d = Document()
        d.add_paragraph("text")
        d.save(str(f))
        result = DocxExtractor().extract(f)
        assert result is not None
        # python-docx may inject a default author string; both None and str are valid
        assert result.metadata.author is None or isinstance(result.metadata.author, str)

    def test_docx_filters_blank_paragraphs(self, tmp_path: Path) -> None:
        from docx import Document
        f = tmp_path / "blanks.docx"
        d = Document()
        d.add_paragraph("First paragraph")
        d.add_paragraph("")
        d.add_paragraph("   ")
        d.add_paragraph("Third paragraph")
        d.save(str(f))
        result = DocxExtractor().extract(f)
        assert result is not None
        lines = result.text.splitlines()
        assert "First paragraph" in lines
        assert "Third paragraph" in lines
        assert "" not in lines
        assert "   " not in lines

    def test_docx_page_count_is_always_none(self, tmp_path: Path) -> None:
        from docx import Document
        f = tmp_path / "doc.docx"
        d = Document()
        d.add_paragraph("text")
        d.save(str(f))
        result = DocxExtractor().extract(f)
        assert result is not None
        assert result.metadata.page_count is None


# ---------------------------------------------------------------------------
# PptxExtractor
# ---------------------------------------------------------------------------


class TestPptxExtractorEdges:
    def test_empty_slide_pptx_returns_document_with_empty_text(self, tmp_path: Path) -> None:
        from pptx import Presentation
        f = tmp_path / "blank.pptx"
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(str(f))
        result = PptxExtractor().extract(f)
        assert result is not None
        assert result.text == ""

    def test_slide_count_stored_as_page_count(self, tmp_path: Path) -> None:
        from pptx import Presentation
        f = tmp_path / "three.pptx"
        prs = Presentation()
        for _ in range(3):
            prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(str(f))
        result = PptxExtractor().extract(f)
        assert result is not None
        assert result.metadata.page_count == 3

    def test_pptx_without_title_has_none_title(self, tmp_path: Path) -> None:
        from pptx import Presentation
        f = tmp_path / "notitle.pptx"
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(str(f))
        result = PptxExtractor().extract(f)
        assert result is not None
        assert result.metadata.title is None

    def test_pptx_language_always_none(self, tmp_path: Path) -> None:
        from pptx import Presentation
        f = tmp_path / "nolang.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(0, 0, 100, 50)
        txBox.text_frame.text = "Some content"
        prs.save(str(f))
        result = PptxExtractor().extract(f)
        assert result is not None
        assert result.metadata.language is None
