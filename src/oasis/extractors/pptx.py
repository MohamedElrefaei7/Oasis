import logging
from pathlib import Path

from pptx import Presentation

from oasis.models import DocumentMetadata, ExtractedDocument

logger = logging.getLogger(__name__)


class PptxExtractor:
    extensions: frozenset[str] = frozenset({".pptx"})

    def extract(self, path: Path) -> ExtractedDocument | None:
        try:
            prs = Presentation(str(path))
        except Exception:
            logger.warning("Failed to open PPTX %s", path, exc_info=True)
            return None

        try:
            lines: list[str] = []
            for slide_idx, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    try:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                lines.append(t)
                    except Exception:
                        logger.debug(
                            "Failed to read shape on slide %d of %s", slide_idx, path, exc_info=True
                        )

            text = "\n".join(lines)
            props = prs.core_properties
            title: str | None = props.title or None
            author: str | None = props.author or None
            slide_count = len(prs.slides)

            stat = path.stat()
            return ExtractedDocument(
                path=path,
                text=text,
                metadata=DocumentMetadata(
                    size_bytes=stat.st_size,
                    mtime=stat.st_mtime,
                    ctime=stat.st_ctime,
                    title=title,
                    author=author,
                    page_count=slide_count,
                ),
            )
        except Exception:
            logger.warning("Failed to extract content from PPTX %s", path, exc_info=True)
            return None
